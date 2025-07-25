import fitz  # PyMuPDF
import docx
import re
import random
from collections import Counter

def parse_file(uploaded_file):
    ext = uploaded_file.name.split('.')[-1]
    
    if ext == 'pdf':
        import fitz  # PyMuPDF
        text = ""
        with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
            for page in doc:
                text += page.get_text()
        return text

    elif ext == 'docx':
        import docx
        doc = docx.Document(uploaded_file)
        return "\n".join([p.text for p in doc.paragraphs])

    elif ext == 'txt':
        return uploaded_file.read().decode('utf-8')

    elif ext == 'pptx':
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    return ""


def summarize_text(text, sentence_count=5):
    sentences = re.split(r'(?<=[.!?]) +', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    top_sentences = sorted(sentences, key=len, reverse=True)[:sentence_count]
    return " ".join(top_sentences)

def brainstorm_ideas(text, top_n=10):
    words = re.findall(r'\b\w{5,}\b', text.lower())
    freq = Counter(words)
    return [w for w, _ in freq.most_common(top_n)]

def extract_keywords(text, top_n=20):
    # Accept shorter keywords (≥3 chars)
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())

    stopwords = set([
        "the", "and", "for", "are", "but", "not", "you", "that", "this", "with", "from", "have", "your",
        "can", "was", "use", "will", "has", "all", "any", "our", "out", "who", "get", "how", "why", "when",
        "had", "what", "where", "which", "into", "about", "also", "other", "more", "than", "then", "them",
        "they", "their", "these", "those", "his", "her", "him", "she", "himself", "herself", "my", "mine"
    ])

    # Filter out common stopwords
    words = [w for w in words if w not in stopwords]
    freq = Counter(words)
    return [w for w, _ in freq.most_common(top_n)]

def generate_quiz(text, num_questions=5):
    import streamlit as st
    import random
    import re

    keywords = extract_keywords(text, top_n=20)
    st.write("📌 Extracted Keywords:", keywords)

    sentences = re.split(r'(?<=[.!?]) +', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
    st.write("📄 Sentence count:", len(sentences))

    quiz = []
    used_keywords = set()

    for kw in keywords:
        if len(quiz) >= num_questions:
            break
        if kw in used_keywords:
            continue

        related_sentences = [s for s in sentences if kw.lower() in s.lower()]
        if not related_sentences:
            continue

        correct = random.choice(related_sentences)
        distractors = random.sample(
            [s for s in sentences if s not in related_sentences],
            k=min(3, len(sentences) - 1)
        )

        question = f"What best describes '{kw}' in the context of this document?"
        options = random.sample([correct] + distractors, len(distractors) + 1)

        quiz.append({
            "question": question,
            "options": options,
            "answer": correct
        })

        used_keywords.add(kw)

    if not quiz:
        quiz.append({
            "question": "⚠️ No quiz could be generated from this file. Try using a file with more full sentences.",
            "options": ["OK", "Retry", "Use PDF", "Exit"],
            "answer": "OK"
        })

    return quiz

def evaluate_answers(quiz, user_answers):
    score = 0
    correct = []
    for i, item in enumerate(quiz):
        if user_answers.get(i) == item["answer"]:
            score += 1
        correct.append(item["answer"])
    return score, correct
